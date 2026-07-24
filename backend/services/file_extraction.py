"""
Extracts plain text from uploaded document bytes, dispatching by file
extension. Each format's extraction logic is isolated in its own helper
so adding a new format later doesn't touch the others.

Deliberately does NOT support legacy binary Office formats (.doc, .ppt) --
only the modern XML-based .docx/.pptx. Also does not do OCR: a
scanned/image-only PDF with no text layer will extract to an empty
string, which raises EmptyExtractionError rather than silently producing
a persona that "knows" an empty document.
"""

import io

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = {"txt", "md", "pdf", "docx", "pptx"}


class UnsupportedFileTypeError(Exception):
    pass


class EmptyExtractionError(Exception):
    """Raised when a file was read successfully but no usable text came
    out of it -- e.g. a scanned PDF with no text layer, or an empty
    file."""


def _extract_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def _extract_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def _extract_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs]
    return "\n".join(paragraphs)


def _extract_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text:
                        lines.append(text)
    return "\n".join(lines)


_EXTRACTORS = {
    "txt": _extract_txt,
    "md": _extract_txt,
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
}


def extract_text(filename: str, content: bytes) -> str:
    """
    Extracts plain text from `content`, choosing the extraction method
    based on `filename`'s extension.

    Raises UnsupportedFileTypeError for extensions not in
    SUPPORTED_EXTENSIONS, and EmptyExtractionError if extraction
    succeeds but produces no usable text (e.g. scanned PDFs).
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFileTypeError(
            f"Unsupported file type '.{ext}'. Supported: "
            f"{', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    text = _EXTRACTORS[ext](content).strip()
    if not text:
        raise EmptyExtractionError(
            f"No extractable text found in '{filename}' -- if this is a "
            "scanned/image-only PDF, it needs OCR first, which isn't "
            "supported yet."
        )
    return text