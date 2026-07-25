import pytest
import io
from docx import Document as DocxDocument
from pptx import Presentation

from backend.services.file_extraction import (
    extract_text,
    UnsupportedFileTypeError,
    EmptyExtractionError,
)


def make_docx_bytes(paragraphs):
    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def make_pptx_bytes(title, body_text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text_frame.text = body_text
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def test_extract_txt():
    text = extract_text("notes.txt", b"Plain text about cell biology.")
    assert text == "Plain text about cell biology."


def test_extract_md():
    text = extract_text("notes.md", b"# Heading\n\nSome markdown content.")
    assert "Heading" in text


def test_extract_docx():
    content = make_docx_bytes(["The mitochondria is the powerhouse of the cell."])
    text = extract_text("notes.docx", content)
    assert "mitochondria" in text


def test_extract_pptx():
    content = make_pptx_bytes("Biology 101", "Cells are the basic unit of life.")
    text = extract_text("slides.pptx", content)
    assert "Biology 101" in text
    assert "basic unit of life" in text


def test_unsupported_extension_raises():
    with pytest.raises(UnsupportedFileTypeError):
        extract_text("notes.ppt", b"legacy binary garbage")


def test_no_extension_raises():
    with pytest.raises(UnsupportedFileTypeError):
        extract_text("noextension", b"some content")


def test_empty_txt_raises_empty_extraction_error():
    with pytest.raises(EmptyExtractionError):
        extract_text("empty.txt", b"   ")


def test_empty_docx_raises_empty_extraction_error():
    content = make_docx_bytes([])
    with pytest.raises(EmptyExtractionError):
        extract_text("empty.docx", content)