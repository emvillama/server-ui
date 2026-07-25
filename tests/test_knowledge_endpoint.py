import io

from docx import Document as DocxDocument
from pptx import Presentation

from backend.services import ollama_client


def make_docx_bytes(paragraphs):
    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def make_pptx_bytes(title, body_text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text_frame.text = body_text
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# --- raw text ingest ---


def test_ingest_knowledge_endpoint(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0, 0.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    resp = client.post(
        f"/personas/{persona['id']}/knowledge",
        json={
            "source_filename": "notes.txt",
            "text": "The mitochondria is the powerhouse of the cell.",
        },
    )
    assert resp.status_code == 201
    body = resp.json()
    assert len(body) == 1
    assert body[0]["source_filename"] == "notes.txt"
    assert body[0]["persona_id"] == persona["id"]
    assert "mitochondria" in body[0]["chunk_text"]
    # embedding vector should NOT be in the response
    assert "embedding" not in body[0]


def test_ingest_knowledge_missing_persona_404s(client, monkeypatch):
    async def fake_embed(model, text):
        return [1.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    resp = client.post(
        "/personas/999999/knowledge",
        json={"source_filename": "notes.txt", "text": "some text"},
    )
    assert resp.status_code == 404


def test_ingest_knowledge_ollama_failure_surfaces_as_502(client, persona, monkeypatch):
    async def failing_embed(model, text):
        raise ollama_client.OllamaError("simulated failure")

    monkeypatch.setattr(ollama_client, "embed", failing_embed)

    resp = client.post(
        f"/personas/{persona['id']}/knowledge",
        json={"source_filename": "notes.txt", "text": "some text"},
    )
    assert resp.status_code == 502


def test_ingest_knowledge_empty_text_returns_422(client, persona):
    resp = client.post(
        f"/personas/{persona['id']}/knowledge",
        json={"source_filename": "notes.txt", "text": ""},
    )
    assert resp.status_code == 422


# --- file upload ---


def test_upload_txt_file(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0, 0.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    resp = client.post(
        f"/personas/{persona['id']}/knowledge/upload",
        files={"file": ("notes.txt", b"Plain text about cell biology.", "text/plain")},
    )
    assert resp.status_code == 201
    body = resp.json()
    assert len(body) == 1
    assert body[0]["source_filename"] == "notes.txt"
    assert "cell biology" in body[0]["chunk_text"]


def test_upload_docx_file(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0, 0.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    content = make_docx_bytes(["The mitochondria is the powerhouse of the cell."])
    resp = client.post(
        f"/personas/{persona['id']}/knowledge/upload",
        files={
            "file": (
                "notes.docx",
                content,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )
    assert resp.status_code == 201
    body = resp.json()
    assert "mitochondria" in body[0]["chunk_text"]


def test_upload_pptx_file(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0, 0.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    content = make_pptx_bytes("Biology 101", "Cells are the basic unit of life.")
    resp = client.post(
        f"/personas/{persona['id']}/knowledge/upload",
        files={
            "file": (
                "slides.pptx",
                content,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert resp.status_code == 201
    body = resp.json()
    combined = " ".join(row["chunk_text"] for row in body)
    assert "Biology 101" in combined
    assert "basic unit of life" in combined


def test_upload_unsupported_extension_returns_422(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    resp = client.post(
        f"/personas/{persona['id']}/knowledge/upload",
        files={"file": ("notes.ppt", b"legacy binary garbage", "application/octet-stream")},
    )
    assert resp.status_code == 422


def test_upload_missing_persona_404s(client, monkeypatch):
    async def fake_embed(model, text):
        return [1.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    resp = client.post(
        "/personas/999999/knowledge/upload",
        files={"file": ("notes.txt", b"some text", "text/plain")},
    )
    assert resp.status_code == 404


# --- overwrite-on-reingest ---


def test_reingesting_same_filename_overwrites_not_duplicates(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    resp1 = client.post(
        f"/personas/{persona['id']}/knowledge",
        json={"source_filename": "notes.txt", "text": "Original content here."},
    )
    assert resp1.status_code == 201

    resp2 = client.post(
        f"/personas/{persona['id']}/knowledge",
        json={"source_filename": "notes.txt", "text": "Updated content here."},
    )
    assert resp2.status_code == 201

    list_resp = client.get(f"/personas/{persona['id']}/knowledge")
    docs = list_resp.json()
    assert len(docs) == 1  # not two -- overwritten, not duplicated
    assert docs[0]["source_filename"] == "notes.txt"
    assert docs[0]["chunk_count"] == 1


# --- list (grouped by document) ---


def test_list_knowledge_groups_by_document(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    client.post(
        f"/personas/{persona['id']}/knowledge",
        json={"source_filename": "a.txt", "text": "First document content."},
    )
    client.post(
        f"/personas/{persona['id']}/knowledge",
        json={"source_filename": "b.txt", "text": "Second document content."},
    )

    resp = client.get(f"/personas/{persona['id']}/knowledge")
    assert resp.status_code == 200
    docs = resp.json()
    assert len(docs) == 2
    filenames = {d["source_filename"] for d in docs}
    assert filenames == {"a.txt", "b.txt"}
    for d in docs:
        assert d["chunk_count"] == 1
        assert "ingested_at" in d


def test_list_knowledge_empty_persona_returns_empty_list(client, persona):
    resp = client.get(f"/personas/{persona['id']}/knowledge")
    assert resp.status_code == 200
    assert resp.json() == []


def test_list_knowledge_missing_persona_404s(client):
    resp = client.get("/personas/999999/knowledge")
    assert resp.status_code == 404


# --- delete ---


def test_delete_document_removes_all_its_chunks(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    client.post(
        f"/personas/{persona['id']}/knowledge",
        json={
            "source_filename": "notes.txt",
            "text": "A " * 400,  # long enough to produce multiple chunks
        },
    )
    list_before = client.get(f"/personas/{persona['id']}/knowledge").json()
    assert len(list_before) == 1

    resp = client.delete(f"/personas/{persona['id']}/knowledge/notes.txt")
    assert resp.status_code == 204

    list_after = client.get(f"/personas/{persona['id']}/knowledge").json()
    assert list_after == []


def test_delete_nonexistent_document_404s(client, persona):
    resp = client.delete(f"/personas/{persona['id']}/knowledge/nope.txt")
    assert resp.status_code == 404


def test_delete_missing_persona_404s(client):
    resp = client.delete("/personas/999999/knowledge/notes.txt")
    assert resp.status_code == 404


def test_delete_only_affects_the_named_document(client, persona, monkeypatch):
    async def fake_embed(model, text):
        return [1.0]

    monkeypatch.setattr(ollama_client, "embed", fake_embed)

    client.post(
        f"/personas/{persona['id']}/knowledge",
        json={"source_filename": "keep.txt", "text": "Keep this content."},
    )
    client.post(
        f"/personas/{persona['id']}/knowledge",
        json={"source_filename": "remove.txt", "text": "Remove this content."},
    )

    resp = client.delete(f"/personas/{persona['id']}/knowledge/remove.txt")
    assert resp.status_code == 204

    remaining = client.get(f"/personas/{persona['id']}/knowledge").json()
    assert len(remaining) == 1
    assert remaining[0]["source_filename"] == "keep.txt"